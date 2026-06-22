import os
import re
from types import SimpleNamespace
from langchain_community.document_loaders import PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
import requests
from bs4 import BeautifulSoup
from sqlalchemy import create_engine, text
from pptx import Presentation

try:
    import pytesseract
    from pdf2image import convert_from_path
except ImportError:
    pytesseract = None
    convert_from_path = None

_DB_URL = None
_engine_cache = {}

def _get_engine(db_url=None):
    url = db_url or _DB_URL or os.getenv("DATABASE_URL", "postgresql+psycopg2://cu:cu@127.0.0.1:5532/cu")
    if url not in _engine_cache:
        _engine_cache[url] = create_engine(url)
    return _engine_cache[url]


def init_collections_table(db_url=None):
    global _DB_URL
    if db_url:
        _DB_URL = db_url
    engine = _get_engine(db_url)
    with engine.begin() as conn:
        conn.execute(text("""
            CREATE TABLE IF NOT EXISTS _collections (
                name VARCHAR(255) PRIMARY KEY
            )
        """))


def load_collections(db_url=None):
    try:
        engine = _get_engine(db_url)
        with engine.connect() as conn:
            result = conn.execute(text("SELECT name FROM _collections ORDER BY name"))
            return [row[0] for row in result]
    except Exception:
        return []


def add_collection(collection_name, db_url=None):
    engine = _get_engine(db_url)
    with engine.begin() as conn:
        conn.execute(
            text("INSERT INTO _collections (name) VALUES (:name) ON CONFLICT (name) DO NOTHING"),
            {"name": collection_name}
        )


def remove_collection(collection_name, db_url=None):
    engine = _get_engine(db_url)
    with engine.begin() as conn:
        conn.execute(text("DELETE FROM _collections WHERE name = :name"), {"name": collection_name})


def validate_collection_name(name):
    return bool(re.match(r'^[a-zA-Z_][a-zA-Z0-9_]*$', name))


# ── Document Readers ──

class SimplePDFReader:
    def __init__(self, chunk=True, chunk_size=500, chunk_overlap=50):
        self.chunk = chunk
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap

    def load(self, path):
        loader = PyPDFLoader(path)
        docs = loader.load()
        if self.chunk:
            splitter = RecursiveCharacterTextSplitter(
                chunk_size=self.chunk_size, chunk_overlap=self.chunk_overlap
            )
            return splitter.split_documents(docs)
        return docs


class OCRPDFReader:
    def __init__(self, chunk=True, chunk_size=500, chunk_overlap=50):
        self.chunk = chunk
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap

    def load(self, path):
        images = convert_from_path(path)
        docs = []
        for i, image in enumerate(images):
            ocr_text = pytesseract.image_to_string(image)
            docs.append(SimpleNamespace(page_content=ocr_text, metadata={"page": i + 1}))
        if self.chunk:
            splitter = RecursiveCharacterTextSplitter(
                chunk_size=self.chunk_size, chunk_overlap=self.chunk_overlap
            )
            return splitter.split_documents(docs)
        return docs


class PPTXReader:
    def __init__(self, chunk=True, chunk_size=500, chunk_overlap=50):
        self.chunk = chunk
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap

    def load(self, path):
        prs = Presentation(path)
        content = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content += shape.text + "\n"
        doc = SimpleNamespace(page_content=content, metadata={"file": path})
        if self.chunk:
            splitter = RecursiveCharacterTextSplitter(
                chunk_size=self.chunk_size, chunk_overlap=self.chunk_overlap
            )
            return splitter.split_documents([doc])
        return [doc]


class WebReader:
    def __init__(self, chunk=True, chunk_size=500, chunk_overlap=50):
        self.chunk = chunk
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap

    def load(self, url):
        try:
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
            }
            response = requests.get(url, headers=headers, timeout=30)
            response.raise_for_status()
            soup = BeautifulSoup(response.content, 'html.parser')
            for script in soup(["script", "style"]):
                script.decompose()
            raw_text = soup.get_text()
            lines = (line.strip() for line in raw_text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            cleaned = ' '.join(chunk for chunk in chunks if chunk)
            doc = SimpleNamespace(
                page_content=cleaned,
                metadata={"url": url, "title": soup.title.string if soup.title else "Unknown"}
            )
            if self.chunk:
                splitter = RecursiveCharacterTextSplitter(
                    chunk_size=self.chunk_size, chunk_overlap=self.chunk_overlap
                )
                return splitter.split_documents([doc])
            return [doc]
        except Exception as e:
            print(f"Error fetching {url}: {e}")
            return []


# ── Vector Store ──

class PgVector2:
    def __init__(self, collection, db_url, embedding_dim=768):
        self.collection = collection
        self.db_url = db_url
        self.embedding_dim = embedding_dim
        self.engine = create_engine(db_url)
        self._ensure_table()

    def _ensure_table(self):
        with self.engine.begin() as conn:
            conn.execute(text("CREATE EXTENSION IF NOT EXISTS vector"))
            conn.execute(text(f'''
                CREATE TABLE IF NOT EXISTS {self.collection} (
                    id SERIAL PRIMARY KEY,
                    content TEXT UNIQUE,
                    embedding VECTOR({self.embedding_dim})
                );
            '''))

    def add_embeddings(self, texts, embeddings):
        with self.engine.begin() as conn:
            for text_val, emb in zip(texts, embeddings):
                emb_str = ','.join(map(str, emb))
                conn.execute(
                    text(f"INSERT INTO {self.collection} (content, embedding) VALUES (:content, :embedding) ON CONFLICT (content) DO NOTHING"),
                    {"content": text_val, "embedding": f'[{emb_str}]'}
                )

    def query(self, embedding, top_k=5):
        emb_str = ','.join(map(str, embedding))
        with self.engine.connect() as conn:
            result = conn.execute(
                text(f"""
                    SELECT content, embedding <#> :embedding AS distance
                    FROM {self.collection}
                    ORDER BY distance ASC
                    LIMIT :k
                """),
                {"embedding": f'[{emb_str}]', "k": top_k}
            )
            return [row[0] for row in result]

    def get_count(self):
        with self.engine.connect() as conn:
            result = conn.execute(text(f"SELECT COUNT(*) FROM {self.collection}"))
            return result.fetchone()[0]

    def drop_table(self):
        with self.engine.begin() as conn:
            conn.execute(text(f"DROP TABLE IF EXISTS {self.collection}"))


# ── Knowledge Bases ──

class PDFKnowledgeBase:
    def __init__(self, path, vector_db, reader):
        self.path = path
        self.vector_db = vector_db
        self.reader = reader

    def load(self, embed_fn):
        docs = self.reader.load(self.path)
        texts = [doc.page_content for doc in docs]
        embeddings = [embed_fn(t) for t in texts]
        self.vector_db.add_embeddings(texts, embeddings)


class WebKnowledgeBase:
    def __init__(self, url, vector_db, reader):
        self.url = url
        self.vector_db = vector_db
        self.reader = reader

    def load(self, embed_fn):
        docs = self.reader.load(self.url)
        texts = [doc.page_content for doc in docs]
        embeddings = [embed_fn(t) for t in texts]
        self.vector_db.add_embeddings(texts, embeddings)


class PPTXKnowledgeBase:
    def __init__(self, path, vector_db, reader):
        self.path = path
        self.vector_db = vector_db
        self.reader = reader

    def load(self, embed_fn):
        docs = self.reader.load(self.path)
        texts = [doc.page_content for doc in docs]
        embeddings = [embed_fn(t) for t in texts]
        self.vector_db.add_embeddings(texts, embeddings)


# ── Multi-Collection Router ──

class MultiKnowledgeBase:
    def __init__(self, db_url, collections, embed_fn=None):
        self.db_url = db_url
        self.collections = collections
        self._embed_fn = embed_fn

    def get_embedding(self, text_input):
        return self._embed_fn(text_input)

    def select_relevant_kb(self, query):
        query_embedding = self.get_embedding(query)
        best_collection = None
        best_score = float('inf')
        for collection in self.collections:
            try:
                vector_db = PgVector2(collection=collection, db_url=self.db_url)
                with vector_db.engine.connect() as conn:
                    result = conn.execute(
                        text(f"""
                            SELECT embedding <#> :embedding AS distance
                            FROM {collection}
                            ORDER BY distance ASC
                            LIMIT 1
                        """),
                        {"embedding": f'[{",".join(map(str, query_embedding))}]'}
                    )
                    row = result.fetchone()
                    if row and row[0] < best_score:
                        best_score = row[0]
                        best_collection = collection
            except Exception as e:
                print(f"Error checking collection {collection}: {e}")
                continue
        return best_collection

    def get_context(self, query, top_k=5):
        best_collection = self.select_relevant_kb(query)
        if not best_collection:
            return []
        print(f"Using knowledge base: {best_collection}")
        vector_db = PgVector2(collection=best_collection, db_url=self.db_url)
        query_embedding = self.get_embedding(query)
        return vector_db.query(query_embedding, top_k=top_k)
